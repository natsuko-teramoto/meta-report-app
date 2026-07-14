from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Inches

from ppt.common.theme import (
    TABLE_HEADER_SIZE,
    TABLE_BODY_SIZE,
    TABLE_HEADER_FILL,
    TABLE_HEADER_TEXT_COLOR,
    TABLE_ALT_ROW_FILL,
    COLOR_TEXT,
    FONT_NAME,
    TABLE_BODY_FILL,
    TABLE_ALT_ROW_FILL,
    TABLE_TOTAL_FILL,
)


def add_metric_table(
    slide,
    df,
    left,
    top,
    width,
    height,
):
    table_data = df.copy()

    rows = len(table_data) + 1
    cols = len(table_data.columns)

    table = slide.shapes.add_table(
        rows,
        cols,
        left,
        top,
        width,
        height,
    ).table

    # 左端の「項目」列だけ少し広くする
    column_weights = [
        1.4 if col_name == "項目" else 1.0
        for col_name in table_data.columns
    ]

    total_weight = sum(column_weights)

    for col_idx, column_weight in enumerate(column_weights):
        table.columns[col_idx].width = int(
            width * column_weight / total_weight
        )

    # 行の高さを表全体に合わせて均等配分
    header_height = Inches(0.32)
    body_height = int(
        (height - header_height) / max(len(table_data), 1)
    )

    table.rows[0].height = header_height

    for row_idx in range(1, rows):
        table.rows[row_idx].height = body_height

    # ヘッダー
    for col_idx, col_name in enumerate(table_data.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)

        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_FILL
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.name = FONT_NAME
        paragraph.font.size = TABLE_HEADER_SIZE
        paragraph.font.bold = True
        paragraph.font.color.rgb = TABLE_HEADER_TEXT_COLOR
        paragraph.alignment = PP_ALIGN.CENTER

    # 本文
    for row_idx, (_, row) in enumerate(
        table_data.iterrows(),
        start=1,
    ):
        row_label = str(row.get("項目", ""))

        for col_idx, col_name in enumerate(table_data.columns):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(row[col_name])
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            cell.fill.solid()

            if row_label.startswith("合計"):
                cell.fill.fore_color.rgb = TABLE_TOTAL_FILL
            elif row_idx % 2 == 0:
                cell.fill.fore_color.rgb = TABLE_ALT_ROW_FILL
            else:
                cell.fill.fore_color.rgb = TABLE_BODY_FILL

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = FONT_NAME
            paragraph.font.size = TABLE_BODY_SIZE
            paragraph.font.color.rgb = COLOR_TEXT

            if col_name == "項目":
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.LEFT
            else:
                paragraph.alignment = PP_ALIGN.RIGHT

            # 合計2行は全セルを太字
            if row_label.startswith("合計"):
                paragraph.font.bold = True

    return table