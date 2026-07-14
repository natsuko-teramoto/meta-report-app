from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from ppt.common.theme import (
    FONT_NAME,
    COLOR_TEXT,
    COLOR_SECTION,
)
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

def add_textbox(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=Pt(12),
    bold=False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.name = FONT_NAME
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = COLOR_TEXT
    p.alignment = PP_ALIGN.LEFT

    return box

def add_section_title(
    slide,
    text,
    left,
    top,
    width,
    height,
):
    # 青い縦棒
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        Pt(4),
        height,
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_SECTION

    accent.line.fill.background()

    # タイトル文字
    return add_textbox(
        slide,
        text,
        left + Pt(10),
        top,
        width - Pt(10),
        height,
        font_size=Pt(14),
        bold=True,
    )