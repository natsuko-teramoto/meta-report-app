from numbers import Number
from pptx.util import Pt
from ppt.common.theme import (
    FONT_NAME,
    COLOR_TEXT,
    CARD_LABEL_SIZE,
    CARD_VALUE_SIZE,
    CARD_DELTA_SIZE,
)

def format_metric_parts(
    value,
    frequency=None,
    reach_rate=None,
):
    if isinstance(value, Number):
        main_text = f"{value:,.0f}"
    else:
        main_text = str(value)

    sub_text = ""

    if frequency is not None:
        sub_text = f" ({frequency:.2f})"

    elif reach_rate is not None:
        sub_text = f" ({reach_rate:.2f}%)"

    return main_text, sub_text

def add_info_card(
    slide,
    label,
    value,
    left,
    top,
    width,
    height,
):
    box = slide.shapes.add_textbox(
        left,
        top,
        width,
        height,
    )

    text_frame = box.text_frame
    text_frame.clear()

    label_paragraph = text_frame.paragraphs[0]
    label_paragraph.text = str(label)
    label_paragraph.font.name = FONT_NAME
    label_paragraph.font.size = CARD_LABEL_SIZE
    label_paragraph.font.color.rgb = COLOR_TEXT

    value_paragraph = text_frame.add_paragraph()
    value_paragraph.text = str(value)
    value_paragraph.font.name = FONT_NAME
    value_paragraph.font.size = CARD_VALUE_SIZE
    value_paragraph.font.bold = True
    value_paragraph.font.color.rgb = COLOR_TEXT

    return box

def add_metric_card(
    slide,
    label,
    value,
    left,
    top,
    width,
    height,
    frequency=None,
    reach_rate=None,
    change=None,
):
    box = slide.shapes.add_textbox(
        left,
        top,
        width,
        height,
    )

    text_frame = box.text_frame
    text_frame.clear()

    label_paragraph = text_frame.paragraphs[0]
    label_paragraph.text = str(label)
    label_paragraph.font.name = FONT_NAME
    label_paragraph.font.size = CARD_LABEL_SIZE
    label_paragraph.font.color.rgb = COLOR_TEXT

    main_text, sub_text = format_metric_parts(
        value,
        frequency,
        reach_rate,
    )

    value_paragraph = text_frame.add_paragraph()

    main_run = value_paragraph.add_run()
    main_run.text = main_text
    main_run.font.name = FONT_NAME
    main_run.font.size = CARD_VALUE_SIZE
    main_run.font.bold = True
    main_run.font.color.rgb = COLOR_TEXT

    if sub_text:
        sub_run = value_paragraph.add_run()
        sub_run.text = sub_text
        sub_run.font.name = FONT_NAME
        sub_run.font.size = Pt(15)
        sub_run.font.bold = False
        sub_run.font.color.rgb = COLOR_TEXT

    if change is not None:
        change_paragraph = text_frame.add_paragraph()

        if change > 0:
            change_paragraph.text = f"前月比 ▲{change:.1f}%"
        elif change < 0:
            change_paragraph.text = f"前月比 ▼{abs(change):.1f}%"
        else:
            change_paragraph.text = "前月比 ±0.0%"

        change_paragraph.font.name = FONT_NAME
        change_paragraph.font.size = CARD_DELTA_SIZE
        change_paragraph.font.color.rgb = COLOR_TEXT

    return box