from pptx import Presentation
from ppt.common.theme import SLIDE_WIDTH, SLIDE_HEIGHT


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])