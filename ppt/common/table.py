from pptx.enum.text import PP_ALIGN

from ppt.common.theme import (
    TABLE_HEADER_SIZE,
    TABLE_BODY_SIZE,
    TABLE_HEADER_FILL,
    TABLE_HEADER_TEXT_COLOR,
    TABLE_BODY_FILL,
    TABLE_ALT_ROW_FILL,
    TABLE_HEADER_HEIGHT,
    TABLE_ROW_HEIGHT,
)

LEFT_ALIGN_COLUMNS = {
    "配置",
    "配信",
    "アトリビューション設定",
    "期間",
}

COLUMN_WIDTH_WEIGHTS = {
    "配置": 1.5,
    "配信": 0.8,
    "アトリビューション設定": 1.9,
    "期間": 2.2,
    "インプレッション": 1.2,
    "リーチ": 1.0,
    "クリック(すべて)": 1.3,
    "リンククリック（すべて）": 1.5,
    "リンククリック": 1.2,
    "LPビュー": 1.0,
}

def add_table(
    slide,
    df,
    left,
    top,
    width,
    height,
    max_rows=4,
):
    table_data = df.head(max_rows)

    rows = len(table_data) + 1
    cols = len(table_data.columns)

    table = slide.shapes.add_table(
        rows,
        cols,
        left,
        top,
        width,
        height,
    ).table

    column_weights = [
        COLUMN_WIDTH_WEIGHTS.get(col_name, 1.0)
        for col_name in table_data.columns
    ]

    total_weight = sum(column_weights)

    for col_idx, column_weight in enumerate(column_weights):
        table.columns[col_idx].width = int(
            width * column_weight / total_weight
        )    

    for col_idx, col_name in enumerate(table_data.columns):
        table.cell(0, col_idx).text = str(col_name)

    for col_idx in range(cols):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_FILL

        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = TABLE_HEADER_SIZE
        p.font.color.rgb = TABLE_HEADER_TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

    table.rows[0].height = TABLE_HEADER_HEIGHT

    for row_idx in range(1, rows):
        table.rows[row_idx].height = TABLE_ROW_HEIGHT

    for row_idx, (_, row) in enumerate(table_data.iterrows(), start=1):
        for col_idx, col_name in enumerate(table_data.columns):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(row[col_name])

            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_BODY_FILL

            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = TABLE_ALT_ROW_FILL

            p = cell.text_frame.paragraphs[0]
            p.font.size = TABLE_BODY_SIZE

            if col_name in LEFT_ALIGN_COLUMNS:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.RIGHT

    return table