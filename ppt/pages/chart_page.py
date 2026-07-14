from pathlib import Path
from pptx.util import Inches

from ppt.common.layout import add_blank_slide
from ppt.common.text import add_textbox
from ppt.common.theme import (
    PAGE_LEFT,
    PAGE_TOP,
    SLIDE_TITLE_SIZE,

    CHART_TITLE_WIDTH,
    CHART_TITLE_HEIGHT,

    CHART_SINGLE_LEFT,
    CHART_SINGLE_TOP,
    CHART_SINGLE_WIDTH,
    CHART_SINGLE_HEIGHT,

    CHART_EMPTY_LEFT,
    CHART_EMPTY_TOP,
    CHART_EMPTY_WIDTH,
    CHART_EMPTY_HEIGHT,

    CHART_MULTI_LAYOUTS,
    CHART_MULTI_DEFAULT_LAYOUT,
)


def add_chart_slide(prs, title, image_path):
    slide = add_blank_slide(prs)

    add_textbox(
        slide,
        title,
        PAGE_LEFT,
        PAGE_TOP,
        CHART_TITLE_WIDTH,
        CHART_TITLE_HEIGHT,
        font_size=SLIDE_TITLE_SIZE,
        bold=True,
    )

    image_path = Path(image_path)

    if image_path.exists():
        slide.shapes.add_picture(
            str(image_path),
            CHART_SINGLE_LEFT,
            CHART_SINGLE_TOP,
            width=CHART_SINGLE_WIDTH,
            height=CHART_SINGLE_HEIGHT,
        )
    else:
        add_textbox(
            slide,
            "グラフ画像がありません",
            CHART_EMPTY_LEFT,
            CHART_EMPTY_TOP,
            CHART_EMPTY_WIDTH,
            CHART_EMPTY_HEIGHT,
        )

    return slide


def add_multi_chart_slide(prs, title, image_paths):
    slide = add_blank_slide(prs)

    add_textbox(
        slide,
        title,
        PAGE_LEFT,
        PAGE_TOP,
        CHART_TITLE_WIDTH,
        CHART_TITLE_HEIGHT,
        font_size=SLIDE_TITLE_SIZE,
        bold=True,
    )

    image_paths = [Path(p) for p in image_paths if Path(p).exists()]
    count = len(image_paths)

    if count == 0:
        add_textbox(
            slide,
            "グラフ画像がありません",
            CHART_EMPTY_LEFT,
            CHART_EMPTY_TOP,
            CHART_EMPTY_WIDTH,
            CHART_EMPTY_HEIGHT,
        )
        return slide

    # 4枚：2×2
    positions = CHART_MULTI_LAYOUTS.get(
        count,
        CHART_MULTI_DEFAULT_LAYOUT,
    )

    for image_path, pos in zip(image_paths, positions):
        left, top, width, height = pos

        slide.shapes.add_picture(
            str(image_path),
            left,
            top,
            width=width,
        )

    return slide