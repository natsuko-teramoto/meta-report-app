from pptx.util import Inches
from pathlib import Path
from ppt.common.layout import add_blank_slide
from ppt.common.text import (
    add_textbox,
    add_section_title,
)
from ppt.common.theme import (
    PAGE_LEFT,
    PAGE_TOP,
    TITLE_SIZE,
    METRIC_TITLE_WIDTH,
    METRIC_TITLE_HEIGHT,
    METRIC_GENDER_LEFT,
    METRIC_AGE_LEFT,
    METRIC_AGE_TOP,
    METRIC_GENDER_TOP,
    METRIC_AGE_WIDTH,
    METRIC_GENDER_WIDTH,
    METRIC_TABLE_LEFT,
    METRIC_TABLE_TOP,
    METRIC_TABLE_WIDTH,
    METRIC_TABLE_HEIGHT,
    METRIC_TABLE_TITLE_TOP,
    METRIC_TABLE_TITLE_WIDTH,
    METRIC_TABLE_TITLE_HEIGHT
)
from ppt.common.metric_table import add_metric_table


def add_metric_analysis_slide(
    prs,
    title,
    analysis,
):
    slide = add_blank_slide(prs)

    add_textbox(
        slide,
        title,
        PAGE_LEFT,
        PAGE_TOP,
        METRIC_TITLE_WIDTH,
        METRIC_TITLE_HEIGHT,
        font_size=TITLE_SIZE,
        bold=True,
    )

    age_image = analysis["age_image"]
    gender_image = analysis["gender_image"]
    table_df = analysis.get("table")

    # 左：年齢×性別の棒グラフ
    slide.shapes.add_picture(
        str(Path(age_image)),
        METRIC_AGE_LEFT,
        METRIC_AGE_TOP,
        width=METRIC_AGE_WIDTH,
    )

    # 右：男女比の円グラフ
    slide.shapes.add_picture(
        str(Path(gender_image)),
        METRIC_GENDER_LEFT,
        METRIC_GENDER_TOP,
        width=METRIC_GENDER_WIDTH,
    )

    add_section_title(
        slide,
        "年齢・性別分析",
        PAGE_LEFT,
        METRIC_TABLE_TITLE_TOP,
        METRIC_TABLE_TITLE_WIDTH,
        METRIC_TABLE_TITLE_HEIGHT,
    )

    # 下：年齢×男女の数・全体割合
# 下：年齢×男女の数・全体割合
    if table_df is not None and not table_df.empty:

        add_section_title(
            slide,
            "年齢・性別分析",
            PAGE_LEFT,
            METRIC_TABLE_TITLE_TOP,
            METRIC_TABLE_TITLE_WIDTH,
            METRIC_TABLE_TITLE_HEIGHT,
        )

        add_metric_table(
            slide,
            table_df,
            METRIC_TABLE_LEFT,
            METRIC_TABLE_TOP,
            METRIC_TABLE_WIDTH,
            METRIC_TABLE_HEIGHT,
        )

    return slide